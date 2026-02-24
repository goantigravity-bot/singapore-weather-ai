"""
generate_slides.py — 从 report.md 内容生成数据准备 PowerPoint 演示文稿

运行: python3 generate_slides.py
输出: docs/product-data-preparation-report/slides.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

logger = logging.getLogger(__name__)

# ── 颜色常量 ──
DARK_BG = "0F172A"
DARK_CARD = "1E293B"
CREAM_BG = "FFF8E1"
ACCENT_BLUE = "38BDF8"
TEXT_DARK = "1A1A2E"
TEXT_LIGHT = "E2E8F0"
TEXT_MUTED = "94A3B8"
GREEN = "4ADE80"
WARN = "FBBF24"
DANGER = "F87171"
BRONZE_BG = "4E342E"
BRONZE_TEXT = "EFEBE9"
SILVER_BG = "37474F"
SILVER_TEXT = "ECEFF1"
GOLD_BG = "E65100"
GOLD_TEXT = "FFF3E0"


def hex_to_rgb(hex_str):
    from pptx.dml.color import RGBColor
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:], 16))


def set_slide_bg(slide, color_hex):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = hex_to_rgb(color)
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, font_size=12, bold=False, color=TEXT_DARK, spacing_after=6, italic=False):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = hex_to_rgb(color)
    p.font.name = "Calibri"
    p.space_after = Pt(spacing_after)
    return p


def add_dark_rect(slide, left, top, width, height, fill_color=DARK_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_table(slide, left, top, width, rows, cols, data,
              header_bg=DARK_CARD, header_fg=ACCENT_BLUE, cell_bg=DARK_BG, cell_fg=TEXT_LIGHT,
              col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(0.3 * rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.name = "Calibri"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = hex_to_rgb(header_fg)
                else:
                    paragraph.font.color.rgb = hex_to_rgb(cell_fg)
            cell_fill = cell.fill
            cell_fill.solid()
            if row_idx == 0:
                cell_fill.fore_color.rgb = hex_to_rgb(header_bg)
            else:
                cell_fill.fore_color.rgb = hex_to_rgb(cell_bg)
            cell.vertical_anchor = MSO_ANCHOR.TOP
    return table_shape


def add_footer(slide, slide_num, is_dark_bg=False):
    """每页页脚：左侧 slide number，中间 classification，右侧留空"""
    text_color = TEXT_MUTED if is_dark_bg else "999999"
    y = 7.1
    # 页码（左）
    add_textbox(slide, 0.6, y, 2, 0.3, f"Slide {slide_num}", 8, color=text_color, align=PP_ALIGN.LEFT)
    # Classification（居中）
    add_textbox(slide, 3, y, 7, 0.3, "Official Closed / Non-Sensitive", 8, color=text_color, align=PP_ALIGN.CENTER)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════════
    # Slide 1: Title (no slide number)
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 1, 1.2, 11, 0.5, "SINGAPORE WEATHER AI", 16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 1.9, 11, 1.2, "Data Preparation Report", 44, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 3.2, 11, 0.5, "From Raw Satellite Data to Production ML Pipeline", 18, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    # 日期和姓名
    add_textbox(slide, 1, 4.5, 11, 0.4, "19-Feb-2026", 14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 4.9, 11, 0.4, "Jin Hui", 14, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 5.8, 11, 0.4, "3-Channel Download In Progress (56.9%)", 11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    # 封面只放 classification，不放 slide number
    add_textbox(slide, 3, 7.1, 7, 0.3, "Official Closed / Non-Sensitive", 8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    slide_num = 1  # 封面不计

    # ═══════════════════════════════════════════
    # Slide 2: Product Vision
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM_BG)
    add_textbox(slide, 0.6, 0.3, 5, 0.4, "PRODUCT VISION", 11, color="999999")
    add_textbox(slide, 0.6, 0.6, 12, 0.8, "Combining satellite cloud patterns with ground sensors to predict rain 10 minutes ahead", 24, bold=True, color=TEXT_DARK)

    add_dark_rect(slide, 0.6, 1.6, 6.5, 5.0)
    metrics = [("10 min", "Forecast Horizon"), ("69", "Ground Stations"), ("3", "Satellite Bands"), ("7", "Sensor Features")]
    for i, (val, label) in enumerate(metrics):
        x = 0.9 + i * 1.55
        add_dark_rect(slide, x, 1.9, 1.35, 0.9, DARK_CARD)
        add_textbox(slide, x, 1.95, 1.35, 0.4, val, 22, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, 2.35, 1.35, 0.4, label, 8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_textbox(slide, 0.9, 3.1, 3, 0.3, "Data Inputs", 12, bold=True, color=ACCENT_BLUE)
    inputs_data = [
        ["Source", "Type", "Data"],
        ["🛰️ Satellite", "Himawari-8/9 L1b", "B08, B11, B13 bands"],
        ["📡 Sensor", "data.gov.sg", "Temp, humidity, rain, PM2.5, wind"],
        ["📍 Location", "GPS embedding", "lat/lon + hour/month cycle"],
    ]
    add_table(slide, 0.9, 3.4, 5.9, 4, 3, inputs_data, col_widths=[1.5, 2.0, 2.4])

    tf = add_rich_textbox(slide, 7.5, 1.6, 5.2, 5.0)
    add_para(tf, "Core idea:", 13, bold=True)
    add_para(tf, "Combine satellite imagery (cloud patterns) with ground sensor readings to predict rainfall at any GPS coordinate in Singapore, 10 minutes ahead.", 11)
    add_para(tf, "The system evolved from a single-channel IR satellite model to a 3-channel multi-sensor fusion architecture over 11 versions in 25 days.", 11)
    add_para(tf, "Key innovation — Local Patch + Coordinate Embedding:", 13, bold=True, spacing_after=2)
    add_para(tf, "Instead of the entire satellite image, the model crops a 32×32 pixel patch directly above each station, learning 'will the cloud overhead produce rain here?'", 11)
    add_para(tf, "Built with React + FastAPI + PyTorch on AWS (EC2 + S3)", 10, color="888888", italic=True)
    add_footer(slide, slide_num)

    # ═══════════════════════════════════════════
    # Slide 3: Version Timeline
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM_BG)
    add_textbox(slide, 0.6, 0.3, 5, 0.4, "DEVELOPMENT JOURNEY", 11, color="999999")
    add_textbox(slide, 0.6, 0.6, 12, 0.8, "Rapid iteration: 11 production releases in 25 days, from idea to multi-channel ML pipeline", 24, bold=True, color=TEXT_DARK)

    versions = [
        ["Version", "Date", "Milestone", "Category"],
        ["v0.1.0", "Jan 26", "Initial full-stack: React + FastAPI, single predict API", "Foundation"],
        ["v0.2–0.4", "Feb 1-4", "Popular search, path query, sensor limit + logging refactor", "Foundation"],
        ["v0.5.0", "Feb 5", "AWS deploy, CloudFront proxy, PM2.5 integration, dashboard", "Infrastructure"],
        ["v0.6.0", "Feb 8", "Training pipeline 7-bug fix, Vitest, multi-recipient email", "Infrastructure"],
        ["v0.7.0", "Feb 10", "SQLite cache layer + ThreadPool parallel inference", "Infrastructure"],
        ["v0.8.0", "Feb 12", "Forecast vs Actual closed-loop verification (SQLite schema)", "Data Quality"],
        ["v0.9.x", "Feb 14-15", "Real-time sensor, WeightedSampler, wind + cloud animation", "Data Quality"],
        ["v0.10.0", "Feb 17", "NOAA data source migration + model tuning experiments", "Optimization"],
        ["v0.11.0", "Feb 17", "Telegram Bot, GCC Terraform, HSD parser (10× perf), 3-ch download", "Optimization"],
    ]
    add_table(slide, 0.6, 1.5, 12, 10, 4, versions,
              header_bg="1A1A2E", header_fg=ACCENT_BLUE, cell_bg="FFFDE7", cell_fg="333333",
              col_widths=[1.5, 1.2, 7.5, 1.8])
    add_footer(slide, slide_num)

    # ═══════════════════════════════════════════
    # Slide 4: Data Source Evolution
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM_BG)
    add_textbox(slide, 0.6, 0.3, 5, 0.4, "SATELLITE DATA", 11, color="999999")
    add_textbox(slide, 0.6, 0.6, 12, 0.8, "3 data source migrations reduced storage from 224TB to 6GB while tripling input channels", 24, bold=True, color=TEXT_DARK)

    add_dark_rect(slide, 0.6, 1.6, 7.2, 5.2)
    add_textbox(slide, 0.9, 1.8, 5, 0.3, "Evolution: JAXA → NOAA → AWS Open Data", 12, bold=True, color=ACCENT_BLUE)

    src_data = [
        ["Aspect", "v1: JAXA FTP", "v2: NOAA ISatSS", "v3: AWS L1b ✅"],
        ["Format", "NetCDF Full Disk", "NetCDF L2 tile", "HSD .bz2 segment"],
        ["Size/file", "~700MB ❌", "~3MB", "~30KB ✅"],
        ["Channels", "C13 (1 ch)", "C13 (1 ch)", "B08+B11+B13 (3 ch) ✅"],
        ["Output", "128×128 .npy", "128×128 .npy", "41×37 .npy × 3 ✅"],
        ["Auth", "JAXA account ❌", "No auth ✅", "No auth ✅"],
        ["6-yr raw", "224 TB ❌", "968 GB", "29 GB ✅"],
        ["6-yr proc", "20 GB", "20 GB", "6 GB ✅"],
    ]
    add_table(slide, 0.9, 2.2, 6.6, 8, 4, src_data, col_widths=[1.2, 1.8, 1.8, 1.8])

    tf = add_rich_textbox(slide, 8.2, 1.6, 4.5, 5.2)
    add_para(tf, "Key decisions:", 13, bold=True)
    add_para(tf, "v1 → v2: JAXA FTP was slow and unreliable across the ocean. Migrated to AWS same-region for speed.", 11)
    add_para(tf, "v2 → v3: 10.1TB of raw NetCDF accumulated on S3 ($53/mo). Switched to L1b HSD — 23,000× smaller per file, 3 channels instead of 1.", 11)
    add_para(tf, "3-channel advantage:", 13, bold=True)
    add_para(tf, "• B08 (6.2μm) — Water vapor / moisture", 11)
    add_para(tf, "• B11 (8.6μm) — Cloud phase (ice vs liquid)", 11)
    add_para(tf, "• B13 (10.4μm) — Cloud top height", 11)
    add_para(tf, "Together, these let the model distinguish rain-producing cumulonimbus from harmless cirrus.", 11)
    add_para(tf, "Expected: +15-25% precision improvement", 11, bold=True, color="1B5E20")
    add_footer(slide, slide_num)

    # ═══════════════════════════════════════════
    # Slide 5: Sensor Features
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM_BG)
    add_textbox(slide, 0.6, 0.3, 5, 0.4, "SENSOR DATA", 11, color="999999")
    add_textbox(slide, 0.6, 0.6, 12, 0.8, "Sensor features grew from 3 to 7 as the model demanded richer signals for tropical rain prediction", 24, bold=True, color=TEXT_DARK)

    add_dark_rect(slide, 0.6, 1.6, 7.2, 5.2)
    add_textbox(slide, 0.9, 1.8, 5, 0.3, "Feature Evolution per Training Phase", 12, bold=True, color=ACCENT_BLUE)

    feat_data = [
        ["Phase", "Features", "Count", "Motivation"],
        ["Phase 1 (v0.1-v0.4)", "🌡️ Temp, 💧 Humidity, 🌧️ Rainfall", "3", "Core weather variables"],
        ["Phase 2 (v0.5)", "+ 🏭 PM2.5", "4", "Haze + NEA outdoor advisory (PSI>100)"],
        ["Phase 3 (v0.9.1)", "+ 🌬️ Wind Speed, 📐 Dir (sin/cos)", "7", "Wind = incoming rain indicator"],
    ]
    add_table(slide, 0.9, 2.2, 6.6, 4, 4, feat_data, col_widths=[1.8, 2.2, 0.6, 2.0])

    add_dark_rect(slide, 0.9, 4.1, 6.4, 0.6, DARK_CARD)
    add_textbox(slide, 1.0, 4.15, 6.2, 0.5, "⚠️ Wind encoding: Direction decomposed to sin(θ) + cos(θ) to avoid circular discontinuity (359° ≈ 1°)", 9, color=WARN)

    add_textbox(slide, 0.9, 5.0, 6, 0.3, "Data Source APIs", 12, bold=True, color=ACCENT_BLUE)
    api_data = [
        ["Endpoint", "Structure", "Update Freq"],
        ["/environment/temperature", "Standard (station-level)", "1 min"],
        ["/environment/humidity", "Standard (station-level)", "1 min"],
        ["/environment/rainfall", "Standard (station-level)", "5 min"],
        ["/environment/pm25", "Regional (mapped to station)", "1 hour"],
        ["/environment/wind-speed + dir", "Standard + sin/cos post-process", "1 min"],
    ]
    add_table(slide, 0.9, 5.3, 6.4, 6, 3, api_data, col_widths=[2.5, 2.5, 1.4])

    tf = add_rich_textbox(slide, 8.2, 1.6, 4.5, 5.2)
    add_para(tf, "Data source: data.gov.sg APIs", 13, bold=True)
    add_para(tf, "5 endpoints returning 3 different JSON structures. Each polled and stored as station-specific CSV files.", 11)
    add_para(tf, "69 weather stations across Singapore provide ground truth.", 11)
    add_para(tf, "PM2.5 motivation:", 13, bold=True)
    add_para(tf, "• Air quality affects visibility; haze correlates with weather patterns", 11)
    add_para(tf, "• Supports outdoor activity recommendations per NEA safety advice (PSI >100 = reduce outdoor exertion)", 11)
    add_para(tf, "Challenges solved:", 13, bold=True)
    add_para(tf, "• PM2.5 hourly vs others per-minute → resample all to 10-min", 11)
    add_para(tf, "• SGT (UTC+8) timestamps vs satellite UTC → conversion", 11)
    add_para(tf, "• PM2.5 regional structure → map to nearest station", 11)
    add_footer(slide, slide_num)

    # ═══════════════════════════════════════════
    # Slide 6: Bronze / Silver / Gold
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM_BG)
    add_textbox(slide, 0.6, 0.3, 5, 0.4, "DATA PIPELINE", 11, color="999999")
    add_textbox(slide, 0.6, 0.6, 12, 0.8, "Data flows through 3 quality tiers — from raw ingestion to training-ready and distributable outputs", 24, bold=True, color=TEXT_DARK)

    # Tier badges — 用与表格一致的术语
    tiers = [
        (1.5, BRONZE_BG, BRONZE_TEXT, "🟫 BRONZE", "Satellite + Sensor"),
        (5.0, SILVER_BG, SILVER_TEXT, "🩶 SILVER", "Satellite + Sensor + Combined"),
        (8.5, GOLD_BG, GOLD_TEXT, "🥇 GOLD", "Training + Frontend"),
    ]
    for x, bg, fg, title, sub in tiers:
        add_dark_rect(slide, x, 1.5, 2.3, 0.8, bg)
        add_textbox(slide, x, 1.5, 2.3, 0.4, title, 13, bold=True, color=fg, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, 1.9, 2.3, 0.3, sub, 9, color=fg, align=PP_ALIGN.CENTER)

    add_textbox(slide, 3.9, 1.6, 0.8, 0.5, "→", 28, bold=True, color="999999", align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.4, 1.6, 0.8, 0.5, "→", 28, bold=True, color="999999", align=PP_ALIGN.CENTER)

    tier_data = [
        ["Tier", "Source", "Location", "Format", "Key Transformation"],
        ["🟫 Bronze", "Satellite", "s3://noaa-himawari8/9/ (public)", "HSD .bz2", "None — raw from AWS Open Data"],
        ["🟫 Bronze", "Sensor", "api.data.gov.sg", "JSON (3 formats)", "None — raw API response"],
        ["🩶 Silver", "Satellite", "s3://…/processed/satellite-3ch/", "41×37 .npy", "Decompress→calibrate→crop→validate→S3"],
        ["🩶 Silver", "Sensor", "s3://…/processed/sensor/ + CSV", "7-feature rows", "Parse→SGT→UTC→resample 10-min→fill gaps"],
        ["🩶 Silver", "Combined", "Time-aligned in S3 processed/", "Paired records", "Match satellite ↔ sensor ±5 min"],
        ["🥇 Gold", "Training", "s3://…/processed/", ".npz per station", "32×32 patch + 6-step LSTM seq + coord"],
        ["🥇 Gold", "Frontend", "API server", "PNG stream", "B13 → grayscale cloud animation overlay"],
    ]
    add_table(slide, 0.6, 2.6, 12, 8, 5, tier_data,
              header_bg="1A1A2E", header_fg=ACCENT_BLUE, cell_bg="FFFDE7", cell_fg="333333",
              col_widths=[1.2, 1.2, 3.5, 2.0, 4.1])
    add_footer(slide, slide_num)

    # ═══════════════════════════════════════════
    # Slide 7: File Availability Issues
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM_BG)
    add_textbox(slide, 0.6, 0.3, 5, 0.4, "DATA CLEANSING", 11, color="999999")
    add_textbox(slide, 0.6, 0.6, 12, 0.8, "Downloading 6 years of satellite data revealed 5 undocumented issues requiring code-level resilience", 24, bold=True, color=TEXT_DARK)

    add_dark_rect(slide, 0.6, 1.6, 7.0, 5.2)
    add_textbox(slide, 0.9, 1.8, 5, 0.3, "⚠️ Issues Not Known Upfront", 12, bold=True, color=WARN)

    issues_data = [
        ["#", "Issue", "Impact", "Code Fix"],
        ["1", "Missing timestamps on S3", "~2-5% slots/day missing", "process_slot() returns 'missing'"],
        ["2", "Himawari-8→9 transition", "Bucket + prefix change at 2022-12-13", "_get_bucket(dt) auto-switch"],
        ["3", "H08 vs H09 filename prefix", "Download fails if wrong prefix", "_s3_key() dynamic prefix"],
        ["4", "NC files use NC_H08_/H09_", "Preprocessing must scan both", "Dual glob pattern scan"],
        ["5", "Sporadic S3 failures", "Single slot should not abort day", "Per-slot try/except"],
    ]
    add_table(slide, 0.9, 2.2, 6.4, 6, 4, issues_data, col_widths=[0.4, 2.0, 2.0, 2.0])

    add_dark_rect(slide, 0.9, 4.8, 6.4, 1.6, DARK_CARD)
    tf = add_rich_textbox(slide, 1.0, 4.9, 6.2, 1.5)
    add_para(tf, "# Himawari-8/9 auto-switching", 9, color=TEXT_MUTED)
    add_para(tf, "H8_END = datetime(2022, 12, 13)", 9, color=WARN)
    add_para(tf, "def _get_bucket(dt):", 9, color=TEXT_LIGHT, spacing_after=2)
    add_para(tf, '    return "noaa-himawari8" if dt < H8_END', 9, color=TEXT_LIGHT, spacing_after=2)
    add_para(tf, '           else "noaa-himawari9"', 9, color=TEXT_LIGHT)

    tf = add_rich_textbox(slide, 8.0, 1.6, 4.8, 5.2)
    add_para(tf, "Key lesson:", 13, bold=True)
    add_para(tf, "When downloading 6 years of historical data across a satellite transition boundary, file naming and bucket locations are not uniform.", 11)
    add_para(tf, "None of these issues are documented in the AWS Open Data catalog.", 11, color="C62828")
    add_para(tf, "The pipeline must be resilient to:", 13, bold=True)
    add_para(tf, "• ~2-5% of timestamps missing per day", 11)
    add_para(tf, "• Bucket name changes mid-dataset", 11)
    add_para(tf, "• Filename prefix changes", 11)
    add_para(tf, "• Network throttling during burst downloads", 11)
    add_para(tf, "Current: 1273/2239 days (56.9%), ETA ~Feb 20. Zero failed days.", 11, bold=True, color="1B5E20")
    add_footer(slide, slide_num)

    # ═══════════════════════════════════════════
    # Slide 8: Per-Station Training
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM_BG)
    add_textbox(slide, 0.6, 0.3, 5, 0.4, "TRAINING STRATEGY", 11, color="999999")
    add_textbox(slide, 0.6, 0.6, 12, 0.8, "Each region gets its own model trained on local satellite + sensor data to capture Singapore's microclimates", 24, bold=True, color=TEXT_DARK)

    add_dark_rect(slide, 0.6, 1.6, 7.0, 5.2)
    add_textbox(slide, 0.9, 1.8, 5, 0.3, "📍 Station Selection per Region", 12, bold=True, color=ACCENT_BLUE)

    station_data = [
        ["Region", "Station", "Why This Site"],
        ["West", "S66 (Tengah)", "Sumatra squall landfall. Open terrain, ~98% completeness"],
        ["Central", "S77 (Marina Bay)", "Busiest pedestrian area — highest user value"],
        ["East", "S107 (East Coast)", "Coastal, NE monsoon exposure. Sea-land boundary"],
        ["North", "S104 (Woodlands)", "Near Johor Strait — Malaysia weather arrives first"],
    ]
    add_table(slide, 0.9, 2.2, 6.4, 5, 3, station_data, col_widths=[1.0, 1.8, 3.6])

    add_dark_rect(slide, 0.9, 4.4, 6.4, 0.5, DARK_CARD)
    add_textbox(slide, 1.0, 4.45, 6.2, 0.4, "Selection: (1) Data completeness  (2) Geographic spread  (3) Pedestrian relevance  (4) Microclimate diversity", 9, color=TEXT_LIGHT)

    add_dark_rect(slide, 0.9, 5.2, 6.4, 1.2, DARK_CARD)
    tf = add_rich_textbox(slide, 1.0, 5.3, 6.2, 1.0)
    add_para(tf, "# Each station generates its own training dataset", 9, color=TEXT_MUTED)
    add_para(tf, "python3 prepare_station_data.py --station S66  # West", 9, color=GREEN)
    add_para(tf, "python3 prepare_station_data.py --station S77  # Central", 9, color=GREEN)
    add_para(tf, "python3 prepare_station_data.py --station S107 # East", 9, color=GREEN)

    tf = add_rich_textbox(slide, 8.0, 1.6, 4.8, 5.2)
    add_para(tf, "Why per-station training?", 13, bold=True)
    add_para(tf, "Singapore is small (~50km) but microclimates differ significantly. A single model averaging all stations loses local patterns.", 11)
    add_para(tf, "Each station's training data pairs its sensor readings with a 32×32 satellite patch directly overhead.", 11)
    add_para(tf, 'The model learns "will it rain at this location" rather than "will it rain somewhere in Singapore."', 11, italic=True, color="666666")
    add_para(tf, "Singapore weather context:", 13, bold=True)
    add_para(tf, "• Convective afternoon rain (2-6 PM daily)", 11)
    add_para(tf, "• NE Monsoon (Dec-Mar): persistent rain", 11)
    add_para(tf, "• SW Monsoon (Jun-Sep): Sumatra squalls from west", 11)
    add_para(tf, "• Rain/dry imbalance: ~13% rain at 1.0mm", 11)
    add_para(tf, "• ENSO cycle: need 6 years to cover", 11)
    add_footer(slide, slide_num)

    # ═══════════════════════════════════════════
    # Slide 9: Model Iterations — with F1 explanation
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM_BG)
    add_textbox(slide, 0.6, 0.3, 5, 0.4, "MODEL TUNING", 11, color="999999")
    add_textbox(slide, 0.6, 0.6, 12, 0.8, "6 architecture experiments to maximize F1 — the balance between catching real rain and avoiding false alarms", 24, bold=True, color=TEXT_DARK)

    add_dark_rect(slide, 0.6, 1.6, 7.0, 2.6)
    add_textbox(slide, 0.9, 1.8, 5, 0.3, "🧪 Architecture Iterations", 12, bold=True, color=ACCENT_BLUE)

    model_data = [
        ["Round", "Change", "F1", "Verdict"],
        ["Baseline", "Conv2d + GAP + LSTM", "55.7%", "Starting point"],
        ["R1", "+ Spatial Attention", "57.7%", "✅ Better cloud focus"],
        ["R2", "+ Deeper CNN (5 layers)", "55.8%", "❌ Overfits"],
        ["R3", "+ Residual Blocks", "54.5%", "❌ Too shallow"],
        ["R4 ✅", "Local Patch + Coord Embed", "61.5%", "✅ BEST"],
        ["R5", "+ Dropout2d in CNN", "—", "❌ Patch too small"],
    ]
    add_table(slide, 0.9, 2.2, 6.4, 7, 4, model_data, col_widths=[1.0, 2.4, 0.8, 2.2])

    add_dark_rect(slide, 0.6, 5.0, 7.0, 1.8)
    add_textbox(slide, 0.9, 5.1, 5, 0.3, "🎯 Rain Threshold Evolution", 12, bold=True, color=ACCENT_BLUE)
    threshold_data = [
        ["Threshold", "Rain %", "F1", "Problem"],
        ["0.1mm/10min", "30.3%", "46.6%", "Too sensitive — drizzle/dew triggers"],
        ["5.0mm/10min", "2.0%", "35.7%", "Too strict — only 2% positive samples"],
        ["1.0mm ✅", "13.0%", "61.5%", "'Need an umbrella' rain ✅"],
    ]
    add_table(slide, 0.9, 5.5, 6.4, 4, 4, threshold_data, col_widths=[1.4, 0.8, 0.6, 3.6])

    # Right side — with F1 explanation
    tf = add_rich_textbox(slide, 8.0, 1.6, 4.8, 5.2)
    add_para(tf, "What is F1 Score?", 14, bold=True)
    add_para(tf, "F1 is the harmonic mean of Precision and Recall — the single best metric for imbalanced classification like rain prediction:", 11)
    add_para(tf, "• Precision: When we say 'it will rain', how often are we right?", 11)
    add_para(tf, "• Recall: Of all actual rain events, how many did we catch?", 11)
    add_para(tf, "• F1 = 2 × (Precision × Recall) / (Precision + Recall)", 11, bold=True, color="0D47A1")
    add_para(tf, "Why F1 matters for this product:", 14, bold=True)
    add_para(tf, "• Low Precision → too many false alarms → users ignore alerts", 11)
    add_para(tf, "• Low Recall → missed rain events → users get caught in rain", 11)
    add_para(tf, "• F1 penalizes both equally — the user experience sweet spot", 11)
    add_para(tf, "R4 breakthrough: Local Patch + Coord Embedding forces the model to focus on the cloud directly overhead + time-of-day context.", 10, italic=True, color="888888")
    add_footer(slide, slide_num)

    # ═══════════════════════════════════════════
    # Slide 10: Infrastructure
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM_BG)
    add_textbox(slide, 0.6, 0.3, 5, 0.4, "INFRASTRUCTURE", 11, color="999999")
    add_textbox(slide, 0.6, 0.6, 12, 0.8, "Minimal viable infrastructure: 3 purpose-built servers at $78/month beat managed platforms for our data scale", 24, bold=True, color=TEXT_DARK)

    add_dark_rect(slide, 0.6, 1.6, 7.0, 2.6)
    add_textbox(slide, 0.9, 1.8, 5, 0.3, "🏗️ Server Topology", 12, bold=True, color=ACCENT_BLUE)
    infra_data = [
        ["Server", "Type", "Role", "Cost"],
        ["API Server", "t3.medium", "FastAPI + React frontend", "~$1/mo (free tier)"],
        ["Training", "g4dn.xlarge", "GPU training (NVIDIA T4)", "~$62/mo (Spot)"],
        ["Download", "t3.xlarge", "Satellite data ingestion", "~$15/mo"],
    ]
    add_table(slide, 0.9, 2.2, 6.4, 4, 4, infra_data, col_widths=[1.4, 1.4, 2.4, 1.2])

    add_dark_rect(slide, 0.6, 4.4, 7.0, 2.4)
    add_textbox(slide, 0.9, 4.6, 5, 0.3, "💰 Cost Optimization", 12, bold=True, color=ACCENT_BLUE)
    cost_data = [
        ["Action", "Before", "After", "Saving"],
        ["S3 raw NC cleanup (10.1TB)", "$53/mo", "<$1/mo", "98%"],
        ["GPU Spot vs On-Demand", "$210/mo", "$62/mo", "70%"],
        ["HSD parser: 3d vs 13d download", "$65", "$15", "77%"],
    ]
    add_table(slide, 0.9, 5.0, 6.4, 4, 4, cost_data, col_widths=[2.4, 1.0, 1.0, 1.0])

    tf = add_rich_textbox(slide, 8.0, 1.6, 4.8, 5.2)
    add_para(tf, "Why not Databricks / Snowflake?", 13, bold=True)
    add_para(tf, "• Data volume (~6GB) too small for cluster overhead", 11)
    add_para(tf, "• Custom binary format (.npy) doesn't fit SQL/tabular", 11)
    add_para(tf, "• Need direct GPU access for custom CNN+LSTM", 11)
    add_para(tf, "• Managed platforms add $100-200/mo without benefit", 11)
    add_para(tf, "Key decisions:", 13, bold=True)
    add_para(tf, "• 3-server split: decouple data-intensive from latency-sensitive", 11)
    add_para(tf, '• S3 as data lake: "Just-in-Time" processing (150GB raw → 5MB patches → purge)', 11)
    add_para(tf, "• Spot Instance for training: 70% cost saving, acceptable for batch ML", 11)
    add_para(tf, "Reconsider managed platforms at >100GB or team size >3", 10, italic=True, color="888888")
    add_footer(slide, slide_num)

    # ═══════════════════════════════════════════
    # Slide 11: Next Steps — with F1 context
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM_BG)
    add_textbox(slide, 0.6, 0.3, 5, 0.4, "NEXT STEPS", 11, color="999999")
    add_textbox(slide, 0.6, 0.6, 12, 0.8, "6 targeted improvements to push F1 from 61.5% to 72-78% — catching more rain while reducing false alarms", 24, bold=True, color=TEXT_DARK)

    next_data = [
        ["#", "Improvement", "Current (R4)", "Target (R5)", "Expected Impact on F1"],
        ["1", "3-channel input", "Single B13 IR", "B08+B11+B13", "+15-25% Precision — distinguish cloud types"],
        ["2", "6× more data", "~395K (1 yr)", "~2.4M (6 yr)", "Full ENSO cycle → better generalization"],
        ["3", "Weighted loss", "MSE (equal)", "3× penalty for missed rain", "↑ Recall: reduce 164 missed rain events (FN)"],
        ["4", "Multi-frame input", "Single frame", "3 frames (t-20,t-10,t)", "Learn cloud movement → better Precision"],
        ["5", "Asymmetric sampling", "50:50 rain:dry", "60:40 + oversample heavy", "↑ Recall for heavy rain (<2% of samples)"],
        ["6", "Chronological split", "Random 80:20", "Train 2020-24, val 2025-26", "Real-world F1: predict future from past"],
    ]
    add_table(slide, 0.6, 1.5, 12, 7, 5, next_data,
              header_bg="1A1A2E", header_fg=ACCENT_BLUE, cell_bg="FFFDE7", cell_fg="333333",
              col_widths=[0.5, 2.0, 2.0, 2.8, 4.7])

    # F1 explanation bar
    add_dark_rect(slide, 1.0, 5.2, 11, 1.5)
    tf = add_rich_textbox(slide, 1.2, 5.3, 10.6, 1.4)
    add_para(tf, "Expected: F1  61.5% (R4)  →  ~72-78% (R5)    |    Based on published CNN+LSTM literature with multi-channel satellite input", 14, bold=True, color=TEXT_LIGHT)
    add_para(tf, "F1 Score = harmonic mean of Precision (are our rain alerts correct?) and Recall (did we catch all the rain?). Higher F1 = fewer false alarms AND fewer missed rain events — directly translating to user trust and safety.", 10, color=TEXT_MUTED)
    add_footer(slide, slide_num)

    # ═══════════════════════════════════════════
    # Slide 12: Thank You
    # ═══════════════════════════════════════════
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 1, 1.2, 11, 0.5, "SINGAPORE WEATHER AI", 16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 1.9, 11, 1.2, "Thank You", 44, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 3.2, 11, 0.5, "Data Preparation Report v1.0", 18, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 4.5, 11, 0.4, "19-Feb-2026  |  Jin Hui", 14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 5.5, 11, 0.5, "📄 Full report: report.html   |   📁 Reference docs: reference/*.html", 12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, slide_num, is_dark_bg=True)

    # Save
    output_path = "docs/product-data-preparation-report/slides.pptx"
    prs.save(output_path)
    logger.info(f"✅ Saved {output_path} ({len(prs.slides)} slides)")
    return output_path


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(message)s")
    path = create_presentation()
    print(f"✅ {path}")
